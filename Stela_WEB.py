import os
import asyncio
import logging
import json
import flet as ft
from flet import audio
from dotenv import load_dotenv
from aiogram import Bot, Dispatcher, types
from aiogram.filters import Command
from aiogram.types import WebAppInfo, InlineKeyboardMarkup, InlineKeyboardButton, FSInputFile
from groq import Groq
from yandex_music import Client
from docx import Document
from pptx import Presentation
from duckduckgo_search import DDGS

# 1. КОНФИГУРАЦИЯ
load_dotenv()
TELEGRAM_TOKEN = os.getenv("TELEGRAM_TOKEN")
GROQ_API_KEY = os.getenv("GROQ_API_KEY")
YANDEX_TOKEN = os.getenv("YANDEX_TOKEN")
MY_CHAT_ID = os.getenv("MY_CHAT_ID") # Твой ID в телеграм для получения файлов
RENDER_URL = os.getenv("RENDER_URL", "https://onrender.com")

logging.basicConfig(level=logging.INFO)

ai_client = Groq(api_key=GROQ_API_KEY) if GROQ_API_KEY else None
y_client = Client(YANDEX_TOKEN).init() if YANDEX_TOKEN else None
bot = Bot(token=TELEGRAM_TOKEN) if TELEGRAM_TOKEN else None
dp = Dispatcher()

# --- ФУНКЦИИ ИНСТРУМЕНТОВ ---

def create_files(mode, title, content):
    fname = f"{title}.docx" if mode == "DOC" else f"{title}.pptx"
    if mode == "DOC":
        doc = Document()
        doc.add_heading(title, 0)
        doc.add_paragraph(content)
        doc.save(fname)
    else:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts)
        slide.shapes.title.text = title
        slide.placeholders.text = content
        prs.save(fname)
    return fname

async def send_to_tg(file_path):
    if bot and MY_CHAT_ID:
        try:
            await bot.send_document(MY_CHAT_ID, FSInputFile(file_path))
            os.remove(file_path) # Чистим место на Render
            return True
        except Exception as e:
            print(f"Ошибка отправки: {e}")
    return False

def web_search(query):
    try:
        with DDGS() as ddgs:
            results = [r['body'] for r in ddgs.text(query, max_results=3)]
            return "\n".join(results)
    except: return "Поиск временно недоступен."

# --- МОЗГ СТЕЛЫ С ПАМЯТЬЮ ---
memory = []

async def get_stela_ai_response(prompt):
    global memory
    system_msg = (
        "Ты Стела, ИИ-ассистент. Твои команды:\n"
        "[MUSIC] Артист - Трек\n"
        "[DOC] Заголовок | Текст\n"
        "[PPT] Заголовок | Текст\n"
        "[SEARCH] Запрос\n"
        "[REMIND] Время | Суть\n"
        "Отвечай кратко, помни контекст. Если создаешь файл, уведомь об отправке в TG."
    )
    
    memory.append({"role": "user", "content": prompt})
    if len(memory) > 10: memory.pop(0) # Храним последние 10 реплик

    try:
        completion = ai_client.chat.completions.create(
            model="llama-3.1-70b-versatile", # Более мощная модель
            messages=[{"role": "system", "content": system_msg}] + memory
        )
        ans = completion.choices.message.content
        memory.append({"role": "assistant", "content": ans})
        return ans
    except Exception as e: return f"Сбой: {e}"

# --- ИНТЕРФЕЙС FLET ---
async def main_flet(page: ft.Page):
    page.title = "Stela Premium OS"
    page.theme_mode = "dark"
    page.bgcolor = "#050505"
    page.padding = 20
    
    audio_player = audio.Audio(src="", autoplay=True)
    page.overlay.append(audio_player)

    status = ft.Text("STELA ONLINE", color="cyan", weight="bold")
    sphere = ft.Container(
        content=ft.Icon(ft.icons.AUTO_AWESOME, size=80, color="cyan200"),
        width=160, height=160, shape="circle",
        gradient=ft.RadialGradient(colors=["#1a237e", "black"]),
        animate_scale=300,
    )
    
    log_column = ft.Column(scroll="auto", height=200, width=350)
    input_f = ft.TextField(label="Команда Стеле...", border_color="cyan", on_submit=lambda e: run_stela(e))

    async def run_stela(e):
        txt = input_f.value
        if not txt: return
        input_f.value = ""
        sphere.scale = 1.3
        log_column.controls.append(ft.Text(f"Вы: {txt}", color="white70"))
        page.update()

        res = await get_stela_ai_response(txt)

        if "[SEARCH]" in res:
            q = res.replace("[SEARCH]", "").strip()
            found = web_search(q)
            res = await get_stela_ai_response(f"Результат поиска: {found}. Сделай вывод.")

        if "[MUSIC]" in res:
            q = res.replace("[MUSIC]", "").strip()
            s = y_client.search(q)
            if s.tracks:
                t = s.tracks.results[0]
                audio_player.src = t.get_download_info(get_direct_links=True)[0].direct_link
                audio_player.play()
                res = f"Играет: {t.title}"
        
        if "[DOC]" in res or "[PPT]" in res:
            m = "DOC" if "[DOC]" in res else "PPT"
            p = res.replace(f"[{m}]", "").split("|")
            fname = create_files(m, p[0].strip(), p[1].strip() if len(p)>1 else "...")
            sent = await send_to_tg(fname)
            res = f"📄 Файл {fname} создан и отправлен вам в Telegram!" if sent else "Файл создан, но не отправлен."

        log_column.controls.append(ft.Text(f"Стела: {res}", color="cyan"))
        sphere.scale = 1.0
        page.update()

    page.add(
        ft.Center(sphere),
        ft.Divider(height=20, color="transparent"),
        ft.Container(log_column, padding=10, bgcolor="#111111", border_radius=10),
        input_f,
        ft.ElevatedButton("ВЫПОЛНИТЬ", on_click=run_stela, width=400, bgcolor="cyan900")
    )

# --- BOT & START ---
@dp.message(Command("start"))
async def cmd_start(msg: types.Message):
    await msg.answer(f"Ваш ID: {msg.from_user.id}\nИспользуйте его в настройках MY_CHAT_ID.")
    markup = InlineKeyboardMarkup(inline_keyboard=[[InlineKeyboardButton(text="STELA OS", web_app=WebAppInfo(url=RENDER_URL))]])
    await msg.answer("Система готова к работе.", reply_markup=markup)

async def main():
    flet_task = ft.app_async(target=main_flet, view=ft.AppView.WEB_BROWSER, host="0.0.0.0", port=int(os.getenv("PORT", 10000)))
    bot_task = dp.start_polling(bot)
    await asyncio.gather(flet_task, bot_task)

if __name__ == "__main__":
    asyncio.run(main())
